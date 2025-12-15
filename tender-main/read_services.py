import os
import re
import csv
import zipfile
from html import unescape


def _read_file_text(path: str, encoding: str = "utf-8") -> str:
    try:
        with open(path, "r", encoding=encoding, errors="ignore") as f:
            return f.read()
    except Exception:
        return ""


# ---------------- DOCX -----------------


def _read_docx_via_zip(path: str) -> str:
    """
    Чтение .docx без внешних библиотек: берём word/document.xml из zip
    и выпиливаем теги.
    """
    try:
        with zipfile.ZipFile(path) as zf:
            with zf.open("word/document.xml") as doc_xml:
                data = doc_xml.read().decode("utf-8", errors="ignore")
    except Exception:
        return ""

    # заменяем теги на пробелы, оставляем текст
    text = re.sub(r"<(.|\n)*?>", " ", data)
    text = unescape(text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def read_docx(path: str) -> str:
    """
    Чтение .docx:

    1) Пытаемся через python-docx, если установлен;
    2) Если нет — через zip+xml (word/document.xml).
    """
    path = os.path.abspath(path)

    # 1. python-docx (если есть)
    try:
        import docx  # type: ignore

        doc = docx.Document(path)
        parts = []

        for p in doc.paragraphs:
            txt = (p.text or "").strip()
            if txt:
                parts.append(txt)

        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))

        text = "\n".join(parts)
        text = re.sub(r"\s+", " ", text)
        return text.strip()
    except Exception:
        # 2. fallback: zip + xml
        return _read_docx_via_zip(path)


# ---------------- DOC (.doc) -----------------


def read_doc(path: str) -> str:
    """
    Старый .doc. Без textract/win32com корректно не прочитаем.
    Поэтому даём простой fallback: пытаемся открыть как текст.
    """
    path = os.path.abspath(path)
    # пробуем как простой текст
    txt = _read_file_text(path, encoding="cp1251")
    if not txt:
        txt = _read_file_text(path, encoding="utf-8")
    return txt.strip()


# ---------------- XLSX / XLS / CSV -----------------


def read_csv(path: str) -> str:
    path = os.path.abspath(path)
    rows = []
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f, delimiter=";")
            for row in reader:
                joined = " | ".join(col.strip() for col in row if col.strip())
                if joined:
                    rows.append(joined)
    except Exception:
        return ""
    return "\n".join(rows)


def read_xlsx(path: str) -> str:
    """
    Примитивное чтение .xlsx: если есть openpyxl — используем, иначе возвращаем пусто.
    """
    path = os.path.abspath(path)
    try:
        import openpyxl  # type: ignore

        wb = openpyxl.load_workbook(path, data_only=True)
        parts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c not in (None, "")]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception:
        return ""


def read_xls(path: str) -> str:
    """
    Примитивный fallback для .xls. Без xlrd читать сложно, поэтому просто возвращаем пусто.
    """
    return ""


# ---------------- PDF -----------------


def read_pdf(path: str) -> str:
    """
    Чтение PDF, если установлен PyPDF2. Иначе — пустая строка.
    """
    path = os.path.abspath(path)
    try:
        import PyPDF2  # type: ignore

        text_parts = []
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                txt = page.extract_text() or ""
                txt = txt.strip()
                if txt:
                    text_parts.append(txt)
        text = "\n".join(text_parts)
        text = re.sub(r"\s+", " ", text)
        return text.strip()
    except Exception:
        return ""


# ---------------- PPTX -----------------


def read_pptx(path: str) -> str:
    """
    Чтение презентаций, если установлен python-pptx.
    """
    path = os.path.abspath(path)
    try:
        import pptx  # type: ignore

        prs = pptx.Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = (shape.text or "").strip()
                    if txt:
                        parts.append(txt)
        text = "\n".join(parts)
        text = re.sub(r"\s+", " ", text)
        return text.strip()
    except Exception:
        return ""


# ---------------- HTML / XML -----------------


def _strip_tags(text: str) -> str:
    text = re.sub(r"<script(.|\n)*?</script>", " ", text, flags=re.IGNORECASE)
    text = re.sub(r"<style(.|\n)*?</style>", " ", text, flags=re.IGNORECASE)
    text = re.sub(r"<(.|\n)*?>", " ", text)
    text = unescape(text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def read_html(path: str) -> str:
    path = os.path.abspath(path)
    txt = _read_file_text(path, encoding="utf-8")
    if not txt:
        txt = _read_file_text(path, encoding="cp1251")
    if not txt:
        return ""
    return _strip_tags(txt)


def read_xml(path: str) -> str:
    path = os.path.abspath(path)
    txt = _read_file_text(path, encoding="utf-8")
    if not txt:
        txt = _read_file_text(path, encoding="cp1251")
    if not txt:
        return ""
    return _strip_tags(txt)
